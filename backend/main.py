import json
import os
import re
import tempfile
import uuid
from pathlib import Path

import fitz  # pymupdf
from docx import Document
from dotenv import load_dotenv
from fastapi import FastAPI, Header, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from groq import Groq
from pydantic import BaseModel
from pptx import Presentation
from supabase import Client, create_client

# ── ENV ───────────────────────────────────────────────────────────────────────
load_dotenv()

GROQ_API_KEY = os.getenv("GROQ_API_KEY")
if not GROQ_API_KEY:
    raise RuntimeError("GROQ_API_KEY is missing in environment variables")

SUPABASE_URL = os.getenv("SUPABASE_URL")
SUPABASE_SERVICE_ROLE_KEY = os.getenv("SUPABASE_SERVICE_ROLE_KEY")
DOCUMENTS_BUCKET = os.getenv("DOCUMENTS_BUCKET", "documents")

if not SUPABASE_URL:
    raise RuntimeError("SUPABASE_URL is missing in environment variables")

if not SUPABASE_SERVICE_ROLE_KEY:
    raise RuntimeError("SUPABASE_SERVICE_ROLE_KEY is missing in environment variables")

supabase_admin: Client = create_client(SUPABASE_URL, SUPABASE_SERVICE_ROLE_KEY)

# ── APP ──────────────────────────────────────────────────────────────────────
app = FastAPI(title="StudyMate AI", version="0.1.0")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_methods=["*"],
    allow_headers=["*"],
)

# ── MODELS ───────────────────────────────────────────────────────────────────
class BucketUploadBody(BaseModel):
    storage_path: str
    filename: str
    content_type: str | None = None

class TranscriptCreateBody(BaseModel):
    transcript_text: str

# ── FILE TYPE MAPS ───────────────────────────────────────────────────────────
EXTENSION_MAP = {
    ".pdf": "pdf",
    ".pptx": "pptx",
    ".docx": "docx",
}

CONTENT_TYPE_MAP = {
    "application/pdf": "pdf",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": "pptx",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": "docx",
}

# ── PARSERS ──────────────────────────────────────────────────────────────────
def _parse_pdf(path: str) -> str:
    doc = fitz.open(path)
    pages = []
    for i, page in enumerate(doc):
        text = page.get_text()
        if text.strip():
            pages.append(f"[Страница {i + 1}]\n{text}")
    return "\n\n".join(pages)


def _parse_pptx(path: str) -> str:
    prs = Presentation(path)
    slides = []

    for i, slide in enumerate(prs.slides):
        texts = [
            para.text.strip()
            for shape in slide.shapes
            if shape.has_text_frame
            for para in shape.text_frame.paragraphs
            if para.text.strip()
        ]

        if texts:
            slides.append(f"[Слайд {i + 1}]\n" + "\n".join(texts))

    return "\n\n".join(slides)


def _parse_docx(path: str) -> str:
    doc = Document(path)
    parts = [p.text.strip() for p in doc.paragraphs if p.text.strip()]

    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(c.text.strip() for c in row.cells if c.text.strip())
            if row_text:
                parts.append(row_text)

    return "\n\n".join(parts)


PARSERS = {
    "pdf": _parse_pdf,
    "pptx": _parse_pptx,
    "docx": _parse_docx,
}


# ── HELPERS ──────────────────────────────────────────────────────────────────
def _clean(text: str) -> str:
    text = re.sub(r"\x00", "", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    text = re.sub(r"[ \t]{2,}", " ", text)
    return text.strip()


def _resolve_doc_type(content_type: str, filename: str) -> str | None:
    if content_type in CONTENT_TYPE_MAP:
        return CONTENT_TYPE_MAP[content_type]
    return EXTENSION_MAP.get(Path(filename).suffix.lower())


def _extract_bearer_token(authorization: str | None) -> str:
    if not authorization or not authorization.startswith("Bearer "):
        raise HTTPException(401, "Missing Authorization bearer token")
    return authorization.removeprefix("Bearer ").strip()


def _require_user(jwt: str):
    try:
        response = supabase_admin.auth.get_user(jwt)
        user = response.user
        if not user:
            raise HTTPException(401, "Invalid access token")
        return user
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(401, f"Unauthorized: {e}")


def _download_from_bucket(storage_path: str) -> bytes:
    try:
        return supabase_admin.storage.from_(DOCUMENTS_BUCKET).download(storage_path)
    except Exception as e:
        raise HTTPException(404, f"File not found in bucket: {e}")


def _parse_file_from_bytes(data: bytes, filename: str, content_type: str) -> tuple[str, str]:
    file_type = _resolve_doc_type(content_type, filename)
    if not file_type:
        raise HTTPException(422, "Supported formats are only PDF, PPTX, DOCX")

    with tempfile.NamedTemporaryFile(suffix=f".{file_type}", delete=False) as tmp:
        tmp.write(data)
        tmp_path = tmp.name

    try:
        text = _clean(PARSERS[file_type](tmp_path))
    except Exception as e:
        raise HTTPException(500, f"Parsing error: {e}")
    finally:
        try:
            os.unlink(tmp_path)
        except OSError:
            pass

    if not text:
        raise HTTPException(422, "Unable to extract text - document is empty or scanned")

    return text, file_type


def _get_document_row(document_id: str):
    try:
        response = (
            supabase_admin.table("documents")
            .select("id, user_id, file_name, file_type, storage_path, created_at")
            .eq("id", document_id)
            .single()
            .execute()
        )
    except Exception as e:
        raise HTTPException(500, f"Database query failed: {e}")

    data = getattr(response, "data", None)
    if not data:
        raise HTTPException(404, "Document not found")

    return data


def _generate_summary(text: str) -> str:
    prompt = f"""You are a study assistant. Analyze the material and create a structured summary:

## Topic
One sentence about what this material is about.

## Key points
5-7 most important points, each 1-2 sentences.

## Summary
A short summary in 3-4 sentences.

## Self-check questions
3 questions based on the material.

Material:
{text[:12000]}"""

    client = Groq(api_key=GROQ_API_KEY)
    response = client.chat.completions.create(
        model="llama-3.3-70b-versatile",
        messages=[{"role": "user", "content": prompt}],
    )
    return response.choices[0].message.content


def _generate_flashcards(text: str) -> list:
    prompt = f"""You are a study assistant. Create 10 flashcards based on the material.

Respond STRICTLY in JSON format, without markdown or explanations:
{{"flashcards": [{{"front": "Question or term", "back": "Answer or definition"}}]}}

Material:
{text[:10000]}"""

    client = Groq(api_key=GROQ_API_KEY)
    response = client.chat.completions.create(
        model="llama-3.3-70b-versatile",
        messages=[{"role": "user", "content": prompt}],
        response_format={"type": "json_object"},
    )
    parsed = json.loads(response.choices[0].message.content)
    return parsed.get("flashcards", [])

def _get_transcript_row(transcript_id: str):
    try:
        response = (
            supabase_admin.table("transcripts")
            .select("id, transcript_text, created_at")
            .eq("id", transcript_id)
            .single()
            .execute()
        )
    except Exception as e:
        raise HTTPException(500, f"Database query failed: {e}")

    data = getattr(response, "data", None)
    if not data:
        raise HTTPException(404, "Transcript not found")

    return data

# ── ROUTES ───────────────────────────────────────────────────────────────────
@app.get("/health")
def health():
    return {"status": "ok"}


@app.post("/api/upload")
async def upload(
        payload: BucketUploadBody,
        authorization: str | None = Header(default=None, alias="Authorization"),
):
    """
    Bucket-first flow:
    - Flutter uploads file to Supabase Storage
    - Flutter sends storage_path + filename + content_type here
    - Backend verifies user token, downloads file from bucket, parses text
    - Backend stores document metadata in Postgres
    """
    jwt = _extract_bearer_token(authorization)
    user = _require_user(jwt)

    if not payload.storage_path.strip():
        raise HTTPException(422, "storage_path is required")

    data = _download_from_bucket(payload.storage_path)
    text, file_type = _parse_file_from_bytes(
        data=data,
        filename=payload.filename,
        content_type=payload.content_type or "",
    )

    document_id = str(uuid.uuid4())

    try:
        supabase_admin.table("documents").insert(
            {
                "id": document_id,
                "user_id": user.id,
                "storage_path": payload.storage_path,
                "file_name": payload.filename,
                "file_type": file_type,
            }
        ).execute()
    except Exception as e:
        raise HTTPException(500, f"Database insert failed: {e}")

    return {
        "document_id": document_id,
        "user_id": user.id,
        "filename": payload.filename,
        "file_type": file_type,
        "char_count": len(text),
        "preview": text[:500],
        "storage_path": payload.storage_path,
    }


@app.post("/api/summary/{document_id}")
async def summary(
        document_id: str,
        authorization: str | None = Header(default=None, alias="Authorization"),
):
    jwt = _extract_bearer_token(authorization)
    user = _require_user(jwt)

    doc = _get_document_row(document_id)
    if doc["user_id"] != user.id:
        raise HTTPException(403, "You do not have access to this document")

    data = _download_from_bucket(doc["storage_path"])
    text, _ = _parse_file_from_bytes(
        data=data,
        filename=doc["file_name"],
        content_type="",
    )

    try:
        summary_text = _generate_summary(text)
    except Exception as e:
        raise HTTPException(500, f"Error generating summary: {e}")

    try:
        supabase_admin.table("summaries").insert(
            {
                "document_id": document_id,
                "summary_text": summary_text,
            }
        ).execute()
    except Exception:
        pass

    return {
        "document_id": document_id,
        "file_name": doc["file_name"],
        "summary": summary_text,
    }


@app.post("/api/flashcards/{document_id}")
async def flashcards(
        document_id: str,
        authorization: str | None = Header(default=None, alias="Authorization"),
):
    jwt = _extract_bearer_token(authorization)
    user = _require_user(jwt)

    doc = _get_document_row(document_id)
    if doc["user_id"] != user.id:
        raise HTTPException(403, "You do not have access to this document")

    data = _download_from_bucket(doc["storage_path"])
    text, _ = _parse_file_from_bytes(
        data=data,
        filename=doc["file_name"],
        content_type="",
    )

    try:
        cards = _generate_flashcards(text)
    except Exception as e:
        raise HTTPException(500, f"Error generating flashcards: {e}")

    rows = []
    for card in cards:
        if not isinstance(card, dict):
            continue

        question = str(card.get("front", "")).strip()
        answer = str(card.get("back", "")).strip()

        if not question and not answer:
            continue

        rows.append(
            {
                "id": str(uuid.uuid4()),
                "document_id": document_id,
                "question": question,
                "answer": answer,
            }
        )

    try:
        if rows:
            supabase_admin.table("flashcards").insert(rows).execute()
    except Exception as e:
        raise HTTPException(500, f"Database insert failed: {e}")

    return {
        "document_id": document_id,
        "file_name": doc["file_name"],
        "count": len(rows),
        "flashcards": rows,
    }

def _transcribe_audio(data: bytes, filename: str) -> str:
    suffix = Path(filename).suffix.lower() or ".m4a"

    with tempfile.NamedTemporaryFile(suffix=suffix, delete=False) as tmp:
        tmp.write(data)
        tmp_path = tmp.name

    try:
        client = Groq(api_key=GROQ_API_KEY)

        with open(tmp_path, "rb") as audio_file:
            response = client.audio.transcriptions.create(
                file=(filename, audio_file.read()),
                model="whisper-large-v3-turbo",
                response_format="json",
            )
    finally:
        try:
            os.unlink(tmp_path)
        except OSError:
            pass

    text = getattr(response, "text", None) or ""
    text = _clean(text)

    if not text:
        raise HTTPException(422, "Unable to transcribe audio")

    return text

@app.post("/api/lecture-summary")
async def lecture_summary(
        payload: BucketUploadBody,
        authorization: str | None = Header(default=None, alias="Authorization"),
):
    jwt = _extract_bearer_token(authorization)
    user = _require_user(jwt)

    if not payload.storage_path.strip():
        raise HTTPException(422, "storage_path is required")

    data = _download_from_bucket(payload.storage_path)

    try:
        transcript = _transcribe_audio(data, payload.filename)
        summary_text = _generate_summary(transcript)
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(500, f"Lecture summary failed: {e}")

    transcript_id = str(uuid.uuid4())
    summary_id = str(uuid.uuid4())

    try:
        supabase_admin.table("transcripts").insert(
            {
                "id": transcript_id,
                "transcript_text": transcript,
            }
        ).execute()

        supabase_admin.table("transcript_summaries").insert(
            {
                "id": summary_id,
                "transcript_id": transcript_id,
                "summary_text": summary_text,
            }
        ).execute()
    except Exception as e:
        raise HTTPException(500, f"Database insert failed: {e}")

    return {
        "user_id": user.id,
        "filename": payload.filename,
        "storage_path": payload.storage_path,
        "transcript_id": transcript_id,
        "transcript": transcript,
        "summary": summary_text,
    }